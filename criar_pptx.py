"""
Gera os arquivos PPTX:
  02_Normas_Regulamentadoras.pptx  — Visão geral das NRs (NR-01, 05, 06, 15, 16, 17, 23, 35)
  04_NR10.pptx                     — NR-10 aprofundada (15 slides)
  05_Extintores.pptx               — Extintores de Incêndio (10 slides)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

W = Inches(13.33)
H = Inches(7.5)
BLUE  = RGBColor(26, 86, 219)
DARK  = RGBColor(17, 24, 39)
MUTED = RGBColor(107, 114, 128)

def new_prs():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs

def add_tb(slide, text, left, top, width, height,
           bold=False, size=20, color=DARK,
           align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb

def add_slide(prs, title, body_lines, notes, is_cover=False):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)

    margin = Inches(0.6)
    content_w = W - 2 * margin

    if is_cover:
        add_tb(slide, title,
               margin, Inches(1.8), content_w, Inches(1.8),
               bold=True, size=40, color=BLUE, align=PP_ALIGN.CENTER)
        y = Inches(3.8)
        for line in body_lines:
            add_tb(slide, line,
                   margin, y, content_w, Inches(0.6),
                   size=22, color=MUTED, align=PP_ALIGN.CENTER)
            y += Inches(0.6)
    else:
        add_tb(slide, title,
               margin, Inches(0.25), content_w, Inches(0.9),
               bold=True, size=26, color=BLUE)
        y = Inches(1.25)
        for line in body_lines:
            add_tb(slide, line,
                   margin, y, content_w, Inches(0.6),
                   size=18, color=DARK)
            y += Inches(0.58)

    slide.notes_slide.notes_text_frame.text = notes
    return slide


# ============================================================
# 02 — NORMAS REGULAMENTADORAS (visão geral)
# ============================================================
def criar_nrs_geral():
    prs = new_prs()
    slides = [
        dict(is_cover=True,
             title="Normas Regulamentadoras\nVisão Geral",
             body=["NR-01 · NR-05 · NR-06 · NR-15 · NR-16 · NR-17 · NR-23 · NR-35",
                   "SENAI · Prof. Marcos José dos Santos · Unidade Cariacica"],
             notes=(
                 "Bem-vindos à aula de Normas Regulamentadoras. "
                 "As NRs são normas obrigatórias expedidas pelo Ministério do Trabalho e Emprego "
                 "que estabelecem requisitos mínimos de segurança e saúde para todos os "
                 "empregadores e empregados regidos pela CLT. "
                 "Vamos estudar as principais normas que você precisa conhecer."
             )),
        dict(title="O que são as NRs e quem deve cumpri-las",
             body=["• Normas obrigatórias expedidas pelo Ministério do Trabalho e Emprego (MTE).",
                   "• Baseadas no Art. 200 da CLT — aplicam-se a toda empresa com empregados CLT.",
                   "• Fiscalizadas pelos Auditores Fiscais do Trabalho — multa, embargo ou interdição.",
                   "• Existem 38 NRs em vigor, cada uma voltada a um risco ou setor específico.",
                   "• NRs Setoriais: NR-18 (construção), NR-20 (combustíveis), NR-31 (agricultura).",
                   "• A base de todas é a NR-01, chamada de Norma Mãe."],
             notes=(
                 "As Normas Regulamentadoras são obrigatórias para todos que possuem empregados "
                 "com registro em carteira. Não são recomendações — são lei. "
                 "O descumprimento sujeita o empregador a multas, embargo da obra ou "
                 "interdição do estabelecimento. "
                 "Atualmente existem 38 NRs em vigor, além de normas setoriais específicas "
                 "para determinados segmentos da economia."
             )),
        dict(title="NR-01 — A Norma Mãe: Gerenciamento de Riscos",
             body=["• Estabelece as diretrizes gerais que regem todas as demais NRs.",
                   "• GRO — Gerenciamento de Riscos Operacionais: identificar, avaliar e controlar riscos.",
                   "• PGR — Programa de Gerenciamento de Riscos: 2 documentos obrigatórios:",
                   "     Inventário de Riscos: mapeia todos os perigos e riscos da empresa.",
                   "     Plano de Ação: define medidas de controle e prazos para cada risco.",
                   "• DIREITO DE RECUSA: trabalhador pode recusar atividade de risco grave e iminente.",
                   "• Hierarquia de controle: Eliminação > Substituição > EPC > Administrativo > EPI."],
             notes=(
                 "A NR-01 é o alicerce de toda a segurança do trabalho no Brasil. "
                 "Ela introduziu o Gerenciamento de Riscos Operacionais, o GRO, e o Programa de "
                 "Gerenciamento de Riscos, o PGR. "
                 "O PGR é composto pelo Inventário de Riscos, que mapeia todos os perigos, "
                 "e pelo Plano de Ação, que define como controlá-los. "
                 "Um direito fundamental introduzido pela NR-01 é o Direito de Recusa: "
                 "o trabalhador pode recusar tarefa que represente risco grave e iminente "
                 "sem sofrer nenhuma penalidade."
             )),
        dict(title="NR-05 — CIPA: Comissão Interna de Prevenção de Acidentes",
             body=["• Obrigatória para empresas com mais de 20 empregados (varia por setor).",
                   "• Composição PARITÁRIA: 50% indicados pelo empregador + 50% eleitos pelos empregados.",
                   "• Mandato: 1 ano com direito a UMA reeleição.",
                   "• Missão: identificar riscos, elaborar o Mapa de Riscos e propor melhorias.",
                   "• Realizar SIPAT — Semana Interna de Prevenção de Acidentes do Trabalho.",
                   "• Membros eleitos têm estabilidade de emprego durante o mandato + 1 ano após."],
             notes=(
                 "A NR-05 cria a Comissão Interna de Prevenção de Acidentes, a CIPA. "
                 "Sua composição é paritária: metade dos membros é indicada pelo empregador "
                 "e metade é eleita pelos próprios empregados. "
                 "O mandato é de um ano, com direito a uma reeleição. "
                 "A CIPA tem como principal missão identificar riscos no ambiente de trabalho, "
                 "elaborar o Mapa de Riscos e propor soluções ao empregador. "
                 "Os membros eleitos gozam de estabilidade de emprego durante o mandato "
                 "e por mais um ano após o seu término."
             )),
        dict(title="NR-06 — EPI: Equipamento de Proteção Individual",
             body=["• EPI é fornecido GRATUITAMENTE pelo empregador — proibido cobrar do trabalhador.",
                   "• Todo EPI deve ter CA (Certificado de Aprovação) emitido pelo MTE.",
                   "• Obrigações do EMPREGADOR: fornecer, treinar o uso, repor, fiscalizar e registrar.",
                   "• Obrigações do EMPREGADO: usar corretamente, conservar, não modificar e comunicar defeitos.",
                   "• Recusa injustificada ao uso de EPI é ato faltoso (advertência, suspensão, demissão).",
                   "• EPI é a ÚLTIMA linha de defesa — o EPC (proteção coletiva) vem primeiro."],
             notes=(
                 "A NR-06 regulamenta os Equipamentos de Proteção Individual. "
                 "O ponto mais importante: o EPI deve ser fornecido gratuitamente pelo empregador. "
                 "Todo EPI deve ter o Certificado de Aprovação, o CA, emitido pelo Ministério do Trabalho, "
                 "que garante que o equipamento foi testado e aprovado para o risco indicado. "
                 "O trabalhador que se recusa a usar EPI sem justificativa comete ato faltoso, "
                 "podendo ser advertido, suspenso ou demitido por justa causa. "
                 "Lembre-se: o EPI é sempre a última linha de defesa — antes dele, "
                 "devem ser adotadas medidas coletivas de proteção."
             )),
        dict(title="NR-15 — Atividades e Operações Insalubres",
             body=["• Insalubridade = exposição a agentes que causam DANO PROGRESSIVO à saúde.",
                   "• Agentes: físicos (ruído, calor, vibração), químicos (poeiras, gases) e biológicos.",
                   "• Graus e adicionais sobre o salário MÍNIMO:",
                   "     Mínimo (Grau 1): 10% do salário mínimo.",
                   "     Médio  (Grau 2): 20% do salário mínimo.",
                   "     Máximo (Grau 3): 40% do salário mínimo.",
                   "• Comprovado por laudo de médico ou engenheiro de segurança do trabalho."],
             notes=(
                 "A NR-15 trata das atividades insalubres, aquelas que expõem o trabalhador "
                 "a agentes que causam dano progressivo à saúde ao longo do tempo. "
                 "Os agentes insalubres podem ser físicos, como ruído e calor, "
                 "químicos, como poeiras e vapores, ou biológicos, como vírus e bactérias. "
                 "Existem três graus de insalubridade: mínimo com adicional de dez por cento, "
                 "médio com vinte por cento, e máximo com quarenta por cento, "
                 "todos calculados sobre o salário mínimo nacional. "
                 "O enquadramento exige laudo técnico de médico ou engenheiro de segurança."
             )),
        dict(title="NR-16 — Atividades e Operações Perigosas",
             body=["• Periculosidade = risco de FATALIDADE IMEDIATA (explosivo, elétrico, radiação, roubos).",
                   "• Diferença fundamental: NR-15 = dano progressivo. NR-16 = risco imediato de morte.",
                   "• Adicional FIXO de 30% sobre o salário BASE contratual (não sobre o mínimo).",
                   "• Atividades: eletricidade, explosivos, inflamáveis, segurança pessoal, motocicletas.",
                   "• Trabalhador não pode receber insalubridade E periculosidade — escolhe o mais vantajoso.",
                   "• Comprovado por laudo pericial de engenheiro de segurança."],
             notes=(
                 "A NR-16 trata das atividades perigosas, que expõem o trabalhador a risco "
                 "imediato de morte ou lesão grave. "
                 "A grande diferença em relação à insalubridade é esta: a insalubridade causa dano "
                 "progressivo ao longo do tempo, enquanto a periculosidade representa risco imediato. "
                 "O adicional de periculosidade é fixo em trinta por cento sobre o salário base contratual. "
                 "Diferente da insalubridade, que incide sobre o salário mínimo. "
                 "O trabalhador não pode acumular os dois adicionais — deve optar pelo mais vantajoso."
             )),
        dict(title="NR-17 — Ergonomia",
             body=["• Objetivo: adaptar as condições de trabalho às características do trabalhador.",
                   "• Previne LER (Lesão por Esforço Repetitivo) e DORT (Distúrbio Osteomuscular).",
                   "• Abrange: mobiliário, equipamentos, condições ambientais e organização do trabalho.",
                   "• Limites de levantamento manual de peso (sem auxílio): 60 kg homem / 20 kg mulher.",
                   "• Análise Ergonômica do Trabalho (AET) obrigatória para atividades de risco.",
                   "• Pausas, postura correta, altura de bancadas e iluminação são reguladas pela NR-17."],
             notes=(
                 "A NR-17 de Ergonomia tem como objetivo adaptar as condições de trabalho "
                 "às características físicas e psicológicas dos trabalhadores. "
                 "Ela previne as chamadas LER e DORT, lesões causadas por movimentos repetitivos "
                 "e posturas inadequadas. "
                 "A norma regulamenta desde a altura das bancadas e cadeiras até o ritmo de trabalho, "
                 "pausas obrigatórias, iluminação e temperatura do ambiente. "
                 "Para atividades de risco ergonômico, é obrigatória a elaboração da "
                 "Análise Ergonômica do Trabalho, a AET."
             )),
        dict(title="NR-23 — Proteção Contra Incêndios",
             body=["• Obrigações do empregador em proteção contra incêndio:",
                   "• Saídas de emergência suficientes, sinalizadas e SEMPRE desobstruídas.",
                   "• Extintores em número, tipo e localização adequados ao risco (NBR 12693).",
                   "• Brigada de Incêndio treinada conforme o porte e o risco da empresa.",
                   "• Plano de Emergência e simulados periódicos obrigatórios.",
                   "• PPCI (Plano de Prevenção e Proteção) aprovado pelo Corpo de Bombeiros.",
                   "• Fiscalização: Corpo de Bombeiros Militar + Ministério do Trabalho e Emprego."],
             notes=(
                 "A NR-23 trata da proteção contra incêndios no ambiente de trabalho. "
                 "O empregador deve garantir saídas de emergência sempre livres e sinalizadas. "
                 "Os extintores devem ser adequados ao tipo de risco de cada área, "
                 "conforme a NBR doze mil seiscentos e noventa e três. "
                 "A brigada de incêndio é formada por funcionários da própria empresa, "
                 "treinados para a resposta imediata a emergências. "
                 "O Plano de Prevenção e Proteção Contra Incêndios, o PPCI, deve ser "
                 "aprovado pelo Corpo de Bombeiros Militar de cada estado."
             )),
        dict(title="NR-35 — Trabalho em Altura",
             body=["• Aplica-se a toda atividade realizada ACIMA DE 2 METROS do nível de referência.",
                   "• Todo trabalhador em altura deve receber treinamento específico (reciclagem a cada 2 anos).",
                   "• APR — Análise Preliminar de Risco: obrigatória antes de cada trabalho em altura.",
                   "• PTS — Permissão de Trabalho em Altura: documento formal de autorização.",
                   "• EPIs obrigatórios: cinto de segurança tipo paraquedista + talabarte + capacete.",
                   "• Ponto de ancoragem deve suportar no mínimo 15 kN (carga estática)."],
             notes=(
                 "A NR-35 regula o trabalho em altura, definido como qualquer atividade "
                 "realizada acima de dois metros do nível de referência onde exista risco de queda. "
                 "Todo trabalhador que atue em altura deve receber treinamento específico, "
                 "com reciclagem obrigatória a cada dois anos. "
                 "Antes de cada trabalho, é obrigatória a Análise Preliminar de Risco, a APR, "
                 "e a emissão da Permissão de Trabalho em Altura, o PTS. "
                 "Os EPIs obrigatórios incluem cinto de segurança tipo paraquedista, "
                 "talabarte e capacete com jugular."
             )),
        dict(title="Hierarquia de Controle de Riscos (NR-01)",
             body=["1. ELIMINAÇÃO — remover completamente o perigo da atividade ou ambiente.",
                   "2. SUBSTITUIÇÃO — trocar o agente/processo perigoso por um menos arriscado.",
                   "3. CONTROLES DE ENGENHARIA (EPC) — enclausuramento, ventilação, proteções coletivas.",
                   "4. CONTROLES ADMINISTRATIVOS — treinamento, rotatividade, sinalização, procedimentos.",
                   "5. EPI — equipamento de proteção individual como ÚLTIMA linha de defesa.",
                   "Regra: sempre aplicar os controles na maior hierarquia possível antes de recorrer ao EPI."],
             notes=(
                 "A hierarquia de controle de riscos é o princípio fundamental da segurança do trabalho. "
                 "Deve-se sempre tentar eliminar o perigo na origem antes de qualquer outra medida. "
                 "Se não for possível eliminar, tenta-se substituir por algo menos perigoso. "
                 "Em seguida, aplicam-se controles de engenharia, como equipamentos de proteção coletiva. "
                 "Depois, controles administrativos como treinamentos e procedimentos. "
                 "O EPI é sempre o último recurso, quando as demais medidas forem insuficientes. "
                 "Esta lógica está na base da NR-01 e deve guiar todas as decisões de segurança."
             )),
        dict(title="Resumo das NRs Abordadas",
             body=["NR-01: Norma Mãe — GRO, PGR (Inventário + Plano de Ação), Direito de Recusa.",
                   "NR-05: CIPA — paritária, mandato de 1 ano, Mapa de Riscos, SIPAT.",
                   "NR-06: EPI — gratuito, CA obrigatório, última linha de defesa.",
                   "NR-15: Insalubridade — dano progressivo, graus 10%/20%/40% sobre salário mínimo.",
                   "NR-16: Periculosidade — risco imediato, adicional fixo 30% sobre salário base.",
                   "NR-17: Ergonomia — adaptar trabalho ao trabalhador, prevenir LER/DORT.",
                   "NR-23: Incêndio — extintores, brigada, PPCI, saídas de emergência.",
                   "NR-35: Altura — acima de 2 m, APR, PTS, cinto paraquedista."],
             notes=(
                 "Vamos revisar o resumo das normas estudadas. "
                 "A NR-01 é a Norma Mãe que rege o gerenciamento de riscos e garante o direito de recusa. "
                 "A NR-05 cria a CIPA com composição paritária e mandato de um ano. "
                 "A NR-06 torna o EPI gratuito e obrigatório como última linha de defesa. "
                 "A NR-15 trata da insalubridade com adicionais sobre o salário mínimo. "
                 "A NR-16 trata da periculosidade com adicional fixo de trinta por cento. "
                 "A NR-17 adapta as condições de trabalho ao trabalhador prevenindo LER e DORT. "
                 "A NR-23 protege contra incêndios com extintores, brigada e PPCI. "
                 "E a NR-35 regula o trabalho em altura acima de dois metros."
             )),
    ]

    for s in slides:
        add_slide(prs, s["title"], s["body"], s["notes"], s.get("is_cover", False))
    path = "02_Normas_Regulamentadoras.pptx"
    prs.save(path)
    print(f"OK: {path} salvo ({len(slides)} slides)")


# ============================================================
# 04 — NR-10 APROFUNDADA (15 slides)
# ============================================================
def criar_nr10():
    prs = new_prs()
    slides = [
        dict(is_cover=True,
             title="NR-10 — Segurança em Instalações\ne Serviços em Eletricidade",
             body=["Ministério do Trabalho e Emprego — Portaria n.º 3.214/78",
                   "SENAI · Prof. Marcos José dos Santos · Unidade Cariacica"],
             notes=(
                 "Bem-vindos à aula aprofundada sobre a NR-10, a norma que regulamenta "
                 "a segurança em instalações e serviços em eletricidade. "
                 "Diferente da visão geral das NRs, aqui vamos detalhar cada aspecto técnico "
                 "desta norma fundamental para quem trabalha com sistemas elétricos."
             )),
        dict(title="O que é a NR-10 e para quem se aplica",
             body=["• Regulamenta segurança e saúde de quem trabalha com instalações elétricas.",
                   "• Aplica-se em TODAS as fases: projeto, execução, operação, manutenção e desmontagem.",
                   "• Obrigatória para QUALQUER trabalhador exposto a risco elétrico — não só eletricistas.",
                   "• Exemplos: pintor que trabalha perto de cabos, pedreiro em obra com rede elétrica.",
                   "• Complementada pela NBR 5410 (baixa tensão) e demais normas ABNT.",
                   "• Descumprimento: multas (NR-28), embargo de obra, interdição do estabelecimento."],
             notes=(
                 "A NR-10 não é exclusiva para eletricistas. Ela se aplica a qualquer trabalhador "
                 "que possa ser exposto a risco elétrico, mesmo que indiretamente. "
                 "Um pintor que trabalha próximo a fiações expostas, ou um pedreiro em obra com "
                 "instalações energizadas, também estão sujeitos a esta norma. "
                 "A NR-10 abrange todas as fases de uma instalação elétrica: "
                 "desde o projeto até a desmontagem, passando pela operação e manutenção."
             )),
        dict(title="Perfis Profissionais — NR-10 item 10.8",
             body=["CAPACITADO: recebeu treinamentos de profissional habilitado ou qualificado.",
                   "   Autorizado para tarefas específicas e limitadas. SEM registro em conselho.",
                   "QUALIFICADO: concluiu curso técnico específico na área elétrica.",
                   "   Possui certificado profissionalizante. SEM necessidade de registro em conselho.",
                   "HABILITADO: além do curso, possui registro no conselho de classe (CREA / CRT).",
                   "   Máxima autonomia — pode assinar laudos, projetos e Análise de Risco.",
                   "PRONTUÁRIO: documento que registra toda a conformidade das instalações. Obrigação do empregador."],
             notes=(
                 "A NR-10 define três perfis profissionais para quem trabalha com eletricidade. "
                 "O capacitado é treinado por um habilitado ou qualificado para tarefas pontuais "
                 "e não precisa de registro em conselho. "
                 "O qualificado concluiu um curso técnico específico e tem certificado profissionalizante, "
                 "mas também não precisa de registro em conselho. "
                 "O habilitado é o nível mais alto: além da formação técnica ou superior, "
                 "possui registro no CREA ou CRT, e pode assinar laudos e projetos. "
                 "O Prontuário de Instalações é o documento que comprova a conformidade de tudo "
                 "e é obrigação do empregador mantê-lo atualizado."
             )),
        dict(title="Prontuário de Instalações Elétricas",
             body=["• Documento obrigatório que comprova a conformidade da instalação com a NR-10.",
                   "• Deve conter: projetos, memoriais descritivos, diagramas unifilares e funcionais.",
                   "• Inclui: laudos, relatórios de inspeção, ART do responsável técnico.",
                   "• Deve ser mantido ATUALIZADO e disponível para fiscalização a qualquer momento.",
                   "• Responsabilidade do EMPREGADOR — não do técnico ou engenheiro contratado.",
                   "• Ausência ou desatualização do Prontuário sujeita o empregador a autuação imediata."],
             notes=(
                 "O Prontuário de Instalações Elétricas é o documento central exigido pela NR-10. "
                 "Ele deve conter todos os projetos elétricos, memoriais descritivos, "
                 "diagramas unifilares e funcionais, laudos técnicos e a Anotação de "
                 "Responsabilidade Técnica do engenheiro responsável. "
                 "É obrigação do empregador mantê-lo sempre atualizado e disponível "
                 "para qualquer fiscalização do Ministério do Trabalho. "
                 "A ausência ou desatualização do prontuário resulta em autuação imediata."
             )),
        dict(title="Sequência de Controle — 5 Passos (LOTO)",
             body=["Obrigatória ANTES de qualquer intervenção em instalação elétrica:",
                   "1. DESENERGIZAÇÃO: desligar todas as fontes (disjuntores, chaves, religadores).",
                   "2. SECCIONAMENTO: isolar fisicamente o trecho do circuito a ser trabalhado.",
                   "3. VERIFICAR AUSÊNCIA DE TENSÃO: instrumento calibrado — NUNCA presuma.",
                   "4. ATERRAMENTO PROVISÓRIO: aterrar e curto-circuitar condutores do trecho.",
                   "5. SINALIZAÇÃO / LOTO: cadeado individual + placa 'Nao ligue — homem trabalhando'.",
                   "Somente apos todos os 5 passos o trabalho pode ser iniciado com seguranca."],
             notes=(
                 "A sequência de controle é o procedimento mais importante da NR-10. "
                 "Antes de tocar em qualquer instalação elétrica, os cinco passos devem ser "
                 "cumpridos na ordem exata: desenergizar, seccionar, verificar ausência de tensão, "
                 "realizar o aterramento provisório e aplicar o sistema LOTO. "
                 "LOTO significa Lockout e Tagout: bloqueio com cadeado individual e etiqueta de aviso. "
                 "O passo mais crítico e frequentemente negligenciado é o terceiro: "
                 "verificar com instrumento calibrado se realmente não há tensão. "
                 "Nunca presuma que está desenergizado — sempre teste."
             )),
        dict(title="LOTO — Lockout / Tagout em Detalhes",
             body=["LOCKOUT (bloqueio físico):",
                   "   Cadeado individual colocado na chave, disjuntor ou válvula de energia.",
                   "   Cada trabalhador usa SEU cadeado — não pode ser removido por outra pessoa.",
                   "TAGOUT (etiqueta de sinalização):",
                   "   Placa de aviso: 'Nao ligue — Homem trabalhando — Nome do responsavel'.",
                   "Se mais de 1 trabalhador: cada um coloca SEU cadeado (vários cadeados, 1 por pessoa).",
                   "Somente o trabalhador que colocou o cadeado pode removê-lo ao finalizar o servico."],
             notes=(
                 "O sistema LOTO é um procedimento de segurança crítico que evita a reenergização "
                 "acidental de um circuito enquanto alguém está trabalhando nele. "
                 "O Lockout é o bloqueio físico com cadeado individual — cada trabalhador usa o seu próprio, "
                 "e ninguém pode remover o cadeado de outra pessoa. "
                 "O Tagout é a etiqueta de aviso com o nome do responsável e a instrução de não ligar. "
                 "Quando múltiplos trabalhadores atuam no mesmo ponto, cada um coloca seu próprio cadeado, "
                 "e o circuito só pode ser reenergizado quando o último cadeado for removido."
             )),
        dict(title="Trabalho Desenergizado × Trabalho Energizado",
             body=["DESENERGIZADO (preferencial — deve ser sempre a primeira opção):",
                   "   Segue os 5 passos. Menor risco. Exige capacitado ou superior.",
                   "ENERGIZADO (exceção — somente quando inviável a desenergização):",
                   "   Exige justificativa técnica formal (por escrito).",
                   "   Análise de Risco (AR) elaborada por habilitado/qualificado.",
                   "   Permissão de Trabalho (PT) emitida antes do início do servico.",
                   "   Executado por HABILITADO + ao menos 1 capacitado presente.",
                   "   EPIs dielétricos obrigatórios (luvas, capacete B, uniforme anti-chama)."],
             notes=(
                 "A NR-10 é clara: o trabalho em instalação energizada é exceção, nunca a regra. "
                 "A desenergização com os cinco passos deve ser sempre a primeira opção. "
                 "O trabalho energizado só é permitido quando a desenergização for comprovadamente "
                 "inviável ou oferecer risco maior que o próprio trabalho. "
                 "Neste caso, é obrigatório elaborar a Análise de Risco por escrito "
                 "e emitir a Permissão de Trabalho antes de iniciar qualquer atividade. "
                 "O serviço deve ser executado por profissional habilitado "
                 "acompanhado de ao menos mais um trabalhador capacitado."
             )),
        dict(title="Análise de Risco (AR) e Permissão de Trabalho (PT)",
             body=["AR — Análise de Risco:",
                   "   Elaborada por habilitado ou qualificado com conhecimento da instalação.",
                   "   Identifica os perigos, avalia a gravidade e define os controles necessários.",
                   "   Específica para cada serviço — não é genérica.",
                   "PT — Permissão de Trabalho:",
                   "   Documento formal que autoriza o início do trabalho em instalação energizada.",
                   "   Deve conter: local, data, serviço, riscos, EPIs, responsáveis e assinaturas.",
                   "   Válida apenas para aquele serviço específico — encerrada ao final."],
             notes=(
                 "A Análise de Risco e a Permissão de Trabalho são documentos distintos e complementares. "
                 "A AR é elaborada antes, identificando os perigos específicos daquela instalação "
                 "e definindo os controles necessários. "
                 "A PT é o documento formal que autoriza o início do trabalho. "
                 "Ela deve conter todos os dados do serviço, os riscos identificados, "
                 "os EPIs requeridos, os responsáveis envolvidos e as assinaturas de todos. "
                 "A PT é válida apenas para aquele serviço específico e deve ser encerrada "
                 "formalmente ao final da atividade."
             )),
        dict(title="Zonas de Trabalho em Instalações Elétricas",
             body=["Zonas definidas em relação às partes ENERGIZADAS da instalacao:",
                   "ZONA DE RISCO (ZR): proximidade tal que o contato acidental e muito provavel.",
                   "   Acesso: somente habilitados com EPIs dielétricos adequados.",
                   "ZONA CONTROLADA (ZC): ao redor da ZR. Acesso restrito e supervisionado.",
                   "   Acesso: capacitados, qualificados ou habilitados com autorização.",
                   "ZONA DE VIZINHANÇA (ZV): além da ZC. Acesso permitido com informação de riscos.",
                   "As distâncias de cada zona variam conforme o nível de tensão (BT, MT ou AT)."],
             notes=(
                 "A NR-10 define zonas de trabalho ao redor de partes energizadas de uma instalação. "
                 "A Zona de Risco é a mais próxima das partes vivas, onde o contato acidental "
                 "é muito provável. Apenas habilitados com EPIs dielétricos podem atuar nela. "
                 "A Zona Controlada fica ao redor da Zona de Risco e o acesso é restrito "
                 "a pessoal autorizado e qualificado. "
                 "A Zona de Vizinhança é a mais afastada, onde pessoas não autorizadas "
                 "podem estar presentes desde que informadas dos riscos. "
                 "As distâncias exatas de cada zona variam conforme o nível de tensão da instalação."
             )),
        dict(title="Riscos Elétricos — Reconhecer para Prevenir",
             body=["CHOQUE ELÉTRICO: passagem de corrente pelo corpo humano.",
                   "   Efeitos: desde formigamento (1 mA) até fibrilação cardíaca (> 50 mA) e morte.",
                   "ARCO ELÉTRICO: descarga elétrica no ar entre condutores — temperatura > 20.000°C.",
                   "   Projeta partículas incandescentes e emite radiação UV intensa.",
                   "QUEIMADURAS: por contato, por arco ou por ignição do vestuário.",
                   "INCÊNDIO: arco ou superaquecimento podem inflamar materiais próximos.",
                   "TENSÃO DE PASSO: diferença de potencial no solo (fio caído) — afastar em passos curtos."],
             notes=(
                 "Conhecer os riscos elétricos é fundamental para preveni-los. "
                 "O choque elétrico ocorre quando a corrente elétrica passa pelo corpo. "
                 "Correntes acima de cinquenta miliamperes podem causar fibrilação ventricular e morte. "
                 "O arco elétrico é ainda mais perigoso: ele atinge temperaturas acima de vinte mil graus "
                 "Celsius, projeta metal fundido e emite radiação ultravioleta intensa. "
                 "A tensão de passo ocorre quando um condutor cai no chão e cria diferença de potencial "
                 "entre os dois pés. Neste caso, aproxime-se em passos curtos ou pule com os dois pés juntos."
             )),
        dict(title="EPIs Obrigatórios para Trabalho Elétrico (NR-06 + NR-10)",
             body=["LUVAS DIELETRICAS: classificadas por classe de tensão (Classe 0 a 4).",
                   "   Inspecionar ANTES de cada uso: enrolar do punho à ponta verificando furos.",
                   "CAPACETE CLASSE B: isolante até 20.000 V. Com jugular obrigatório.",
                   "OCULOS DE PROTEÇÃO: contra arco elétrico e projeção de partículas.",
                   "CALCADO DIELETRICO: solado isolante, sem partes metálicas expostas.",
                   "UNIFORME ANTI-CHAMA (FR): tecido tratado — proibido qualquer sintético.",
                   "MANGA PROTETORA: usada com a luva dielétrica para proteção do antebraço."],
             notes=(
                 "Os EPIs para trabalho elétrico são determinados conjuntamente pela NR-06 e NR-10. "
                 "As luvas dielétricas são o EPI mais importante: devem ser inspecionadas antes de cada uso "
                 "enrolando-as do punho à ponta para detectar qualquer furo invisível. "
                 "O capacete deve ser da classe B, que isola até vinte mil volts, e deve ter jugular. "
                 "O uniforme anti-chama é feito de tecido tratado que não derrete nem pega fogo facilmente. "
                 "Qualquer tecido sintético como poliéster é proibido, pois derrete e gruda na pele "
                 "em contato com o arco elétrico."
             )),
        dict(title="Ferramentas Dielétricas",
             body=["• Ferramentas com cabo e/ou corpo isolado para uso em instalações elétricas.",
                   "• Classificadas por tensão máxima de trabalho: 1.000 V CA / 1.500 V CC.",
                   "• Regulamentadas pela NBR IEC 60900.",
                   "• Inspeção obrigatória antes de cada uso: verificar trincas, desgaste e integridade.",
                   "• Validade de teste dielétrico: respeitar o prazo impresso na ferramenta.",
                   "• Ferramentas com isolamento danificado devem ser RETIRADAS DE USO imediatamente.",
                   "• Nunca use ferramentas convencionais (não isoladas) em instalações energizadas."],
             notes=(
                 "As ferramentas dielétricas são especialmente projetadas para uso em instalações elétricas. "
                 "Elas possuem cabo e corpo com isolamento que suportam até mil volts em corrente alternada "
                 "ou mil e quinhentos volts em corrente contínua, conforme a NBR IEC 60900. "
                 "Antes de cada uso, o trabalhador deve inspecionar visualmente a ferramenta, "
                 "verificando se há trincas, desgastes ou danos no isolamento. "
                 "Qualquer ferramenta com isolamento comprometido deve ser retirada de uso imediatamente. "
                 "Usar uma ferramenta convencional em instalação energizada é uma violação grave da NR-10."
             )),
        dict(title="Responsabilidades do Empregador (NR-10)",
             body=["• Garantir que as instalações elétricas sejam seguras e estejam em conformidade.",
                   "• Manter o Prontuário de Instalações atualizado e disponível.",
                   "• Fornecer EPIs dielétricos gratuitamente e em quantidade suficiente.",
                   "• Promover treinamento NR-10 a CADA 2 ANOS ou quando houver mudança de risco.",
                   "• Garantir que trabalho energizado tenha AR, PT e habilitado responsável.",
                   "• Implementar e manter o sistema LOTO na empresa.",
                   "• Manter registros de todos os treinamentos, inspeções e manutenções."],
             notes=(
                 "O empregador tem responsabilidades amplas e específicas definidas pela NR-10. "
                 "Ele deve garantir que as instalações sejam seguras e conformes com a norma, "
                 "manter o Prontuário de Instalações sempre atualizado e fornecer gratuitamente "
                 "todos os EPIs dielétricos necessários. "
                 "A reciclagem do treinamento NR-10 é obrigatória a cada dois anos "
                 "ou sempre que houver mudança nas condições ou nos riscos da instalação. "
                 "O sistema LOTO deve ser implementado e mantido como procedimento padrão da empresa."
             )),
        dict(title="Responsabilidades do Trabalhador + Direito de Recusa",
             body=["• Seguir rigorosamente a sequência de controle (5 passos) antes de qualquer intervenção.",
                   "• Usar corretamente TODOS os EPIs fornecidos — nunca improvise.",
                   "• NUNCA presumir que a instalação está desenergizada — teste sempre.",
                   "• Comunicar imediatamente ao superior qualquer condição insegura identificada.",
                   "• Participar de todos os treinamentos e reciclagens promovidos pela empresa.",
                   "• DIREITO DE RECUSA: recusar tarefa de risco grave e iminente SEM PENALIDADES.",
                   "   Comunicar formalmente ao superior e registrar a recusa por escrito."],
             notes=(
                 "O trabalhador também tem responsabilidades claras na NR-10. "
                 "A principal delas é seguir sempre a sequência de controle antes de qualquer intervenção. "
                 "Nunca presuma que uma instalação está desenergizada — sempre teste com instrumento. "
                 "O Direito de Recusa é um direito fundamental: o trabalhador pode e deve recusar "
                 "qualquer tarefa que represente risco grave e iminente à sua integridade. "
                 "Esta recusa deve ser comunicada formalmente ao superior e registrada por escrito. "
                 "O trabalhador não pode sofrer nenhuma penalidade por exercer este direito."
             )),
        dict(title="Resumo — 10 Pontos Essenciais da NR-10",
             body=["1. Aplica-se a TODOS expostos a risco elétrico — não so eletricistas.",
                   "2. Tres perfis: Capacitado (treinado) > Qualificado (certificado) > Habilitado (CREA).",
                   "3. Prontuario obrigatorio — responsabilidade do empregador.",
                   "4. 5 passos LOTO antes de qualquer intervencao elétrica.",
                   "5. NUNCA presuma tensão zero — verifique sempre com instrumento.",
                   "6. Trabalho energizado e EXCECAO — exige AR, PT e habilitado.",
                   "7. EPIs: luvas dieletricas, capacete B, uniforme FR, calcado dieletrico.",
                   "8. Ferramentas: isoladas, dentro da validade, inspecionadas antes do uso.",
                   "9. Reciclagem obrigatoria a cada 2 anos.",
                   "10. Direito de Recusa: recuse risco grave sem medo de penalidades."],
             notes=(
                 "Vamos revisar os dez pontos essenciais da NR-10. "
                 "Primeiro: a norma se aplica a todos expostos ao risco elétrico, não só eletricistas. "
                 "Segundo: existem três perfis profissionais — capacitado, qualificado e habilitado. "
                 "Terceiro: o Prontuário é obrigatório e responsabilidade do empregador. "
                 "Quarto: os cinco passos LOTO devem ser seguidos antes de qualquer intervenção. "
                 "Quinto: nunca presuma tensão zero — sempre verifique com instrumento calibrado. "
                 "Sexto: trabalho energizado é exceção e exige AR, PT e profissional habilitado. "
                 "Sétimo: EPIs dielétricos são obrigatórios e devem ser inspecionados. "
                 "Oitavo: ferramentas devem ser dielétricas e dentro do prazo de validade. "
                 "Nono: reciclagem obrigatória a cada dois anos. "
                 "Décimo: exercite o Direito de Recusa quando houver risco grave e iminente."
             )),
    ]

    for s in slides:
        add_slide(prs, s["title"], s["body"], s["notes"], s.get("is_cover", False))
    path = "04_NR10.pptx"
    prs.save(path)
    print(f"OK: {path} salvo ({len(slides)} slides)")


# ============================================================
# 05 — EXTINTORES DE INCÊNDIO (10 slides)
# ============================================================
def criar_extintores():
    prs = new_prs()
    slides = [
        dict(is_cover=True,
             title="Extintores de Incêndio\nPrevenção e Combate",
             body=["NR-23 · ABNT NBR 12693 · SENAI Unidade Cariacica",
                   "Prof. Marcos José dos Santos"],
             notes=(
                 "Bem-vindos à aula sobre extintores de incêndio. "
                 "Vamos aprender sobre os tipos de incêndio, as classes de fogo, "
                 "os agentes extintores corretos para cada situação e como usar "
                 "um extintor de forma segura com a técnica PASS."
             )),
        dict(title="O Triângulo do Fogo — Como o Incêndio se Forma",
             body=["O fogo precisa de TRES elementos simultâneos:",
                   "  COMBUSTIVEL — material que queima (madeira, papel, gasolina, gás).",
                   "  COMBURENTE — oxigênio do ar (mínimo ~16% para sustentar chama).",
                   "  CALOR (energia de ativação) — fonte de ignição (faísca, chama, atrito).",
                   "O Tetraedro do Fogo acrescenta um 4° elemento: REACAO EM CADEIA.",
                   "Remover QUALQUER elemento extingue o incêndio.",
                   "Principio da extincao: resfriar, abafar, isolar ou interromper a reacao."],
             notes=(
                 "O fogo resulta da combinação de três elementos: combustível, comburente e calor. "
                 "Para apagá-lo, basta eliminar um desses elementos. "
                 "O conceito moderno inclui um quarto elemento, a reação em cadeia, "
                 "formando o Tetraedro do Fogo. "
                 "Os extintores podem agir por resfriamento, abafamento, isolamento do combustível "
                 "ou interrupção da reação em cadeia, dependendo do tipo de agente extintor."
             )),
        dict(title="Classes de Incêndio",
             body=["CLASSE A — Solidos com brasa residual: madeira, papel, tecido, borracha.",
                   "   Agente: AGUA ou po ABC. Extincao por resfriamento.",
                   "CLASSE B — Líquidos e gases inflamaveis: gasolina, alcool, GLP.",
                   "   Agente: PO BC ou ABC, CO2, espuma. NUNCA use agua.",
                   "CLASSE C — Equipamentos elétricos ENERGIZADOS.",
                   "   Agente: PO BC, ABC ou CO2 (nao condutores). NUNCA use agua.",
                   "CLASSE D — Metais combustíveis: magnesio, titanio, sodio. Po especial.",
                   "CLASSE K — Oleos e gorduras de cozinha. Extintor especifico — NUNCA agua."],
             notes=(
                 "Os incêndios são classificados em cinco classes principais. "
                 "Classe A são os sólidos que formam brasa: use água ou pó ABC. "
                 "Classe B são líquidos e gases inflamáveis: nunca use água, que pode espirrar o líquido em chamas. "
                 "Classe C são equipamentos elétricos energizados: nunca use água por risco de choque. "
                 "Classe D são metais combustíveis raros que exigem pó especial específico. "
                 "Classe K são óleos e gorduras de cozinha em alta temperatura: "
                 "nunca use água, pois o vapor explosivo pode projetar óleo incandescente."
             )),
        dict(title="Tipos de Extintores e Agentes",
             body=["AGUA PRESSURIZADA: Classe A. Age por resfriamento. Nao usar em B, C ou D.",
                   "PO QUIMICO SECO BC: Classes B e C. Interrompe reacao em cadeia.",
                   "PO QUIMICO SECO ABC: Classes A, B e C. Mais versatil — uso mais comum.",
                   "GAS CARBONICO (CO2): Classes B e C. Nao deixa residuo. Ideal para TI.",
                   "ESPUMA MECANICA: Classes A e B. Age por abafamento e resfriamento.",
                   "PO ESPECIAL: Classe D. Especifico para cada metal combustivel.",
                   "Atencao: po ABC deixa residuo — pode danar equipamentos electronicos."],
             notes=(
                 "Os principais tipos de extintores diferem pelo agente extintor que utilizam. "
                 "O de água age por resfriamento e serve apenas para a Classe A. "
                 "O pó BC combate as classes B e C interrompendo a reação em cadeia. "
                 "O pó ABC é o mais versátil, cobrindo as três classes mais comuns. "
                 "O CO2 é ideal para equipamentos eletrônicos e ambientes de TI pois não deixa resíduo, "
                 "mas atenção: em espaços fechados pode causar asfixia. "
                 "A espuma é usada para grandes volumes de líquidos inflamáveis."
             )),
        dict(title="Seleção Correta do Extintor",
             body=["Lixeira com papel em chamas (Classe A)      → Agua ou Po ABC.",
                   "Vazamento de GLP inflamado (Classe B)       → Po BC, ABC ou CO2.",
                   "Curto em quadro eletrico energizado (Cl. C) → Po BC, ABC ou CO2. NUNCA agua.",
                   "Servidores e computadores                   → CO2 (sem residuo).",
                   "Oleo de cozinha em alta temperatura (Cl. K) → Extintor especifico po umido.",
                   "Metal combustivel (Classe D)                → Po metalico especifico.",
                   "REGRA DE OURO: leia o rotulo do extintor — ele indica as classes adequadas."],
             notes=(
                 "A escolha correta do extintor é fundamental para a eficácia e para a segurança do operador. "
                 "Para sólidos com brasa, como papel e madeira, use água ou pó ABC. "
                 "Para líquidos e gases inflamáveis, use pó BC, ABC ou CO2, nunca água. "
                 "Para equipamentos elétricos energizados, use pó seco ou CO2, nunca água. "
                 "Para servidores e computadores, prefira CO2 por não deixar resíduo de pó. "
                 "Para óleos de cozinha quentes, use apenas o extintor específico de pó úmido para Classe K. "
                 "Sempre leia o rótulo do extintor antes de usá-lo."
             )),
        dict(title="Como Usar um Extintor — Técnica PASS",
             body=["P — PUXE o pino de segurança (remove o lacre).",
                   "A — APONTE o bico/mangueira para a BASE das chamas, nao para o topo.",
                   "S — APERTE (Squeeze) o gatilho com firmeza e de forma continua.",
                   "S — VARRA (Sweep) da esquerda para a direita cobrindo toda a base.",
                   "Distancia segura: entre 1,5 m e 3 m da chama (consulte o rotulo).",
                   "Posicione-se com o VENTO NAS COSTAS — evita que o agente volte em voce.",
                   "Mantenha SEMPRE uma rota de fuga livre atras de voce."],
             notes=(
                 "A técnica PASS é o método padrão internacional para uso de extintores. "
                 "P de Puxar o pino de segurança que impede o acionamento acidental. "
                 "A de Apontar o bico para a base das chamas, nunca para o topo, "
                 "pois é na base que está o combustível. "
                 "S de Apertar o gatilho com firmeza e de forma contínua. "
                 "S de Varrer o jato da esquerda para a direita cobrindo toda a base. "
                 "Posicione-se sempre com o vento nas costas e mantenha uma rota de fuga livre."
             )),
        dict(title="Inspeção, Manutenção e Validade",
             body=["INSPECAO MENSAL (visual) — pelo proprio funcionário treinado:",
                   "   Lacre intacto, manometro na zona VERDE, sem danos ou corrosao.",
                   "   Local acessível, sinalizado e desobstruido.",
                   "MANUTENCAO ANUAL — empresa credenciada pelo INMETRO.",
                   "   Verificacao interna, recarga, teste do mecanismo e laudo técnico.",
                   "TESTE HIDROSTATICO — a cada 5 anos para cilindros de aco.",
                   "APOS QUALQUER USO (mesmo parcial) — recarregar imediatamente.",
                   "Extintor com manometro na zona vermelha = INOPERANTE."],
             notes=(
                 "A manutenção regular dos extintores é tão importante quanto saber usá-los. "
                 "A inspeção visual mensal deve ser feita por funcionário treinado, verificando "
                 "o lacre, a pressão no manômetro, que deve estar na zona verde, "
                 "e se o extintor está acessível e sinalizado. "
                 "A manutenção anual deve ser realizada por empresa credenciada pelo INMETRO, "
                 "que realiza verificação interna, recarga e emite laudo técnico. "
                 "O teste hidrostático do cilindro é obrigatório a cada cinco anos. "
                 "Após qualquer uso, mesmo que parcial, o extintor deve ser recarregado imediatamente."
             )),
        dict(title="Localização e Sinalização (NR-23 / NBR 12693)",
             body=["ALTURA: parte inferior do extintor a no maximo 1,60 m do piso.",
                   "DISTANCIA MAXIMA a percorrer:",
                   "   Classe A: ate 15 metros.",
                   "   Classes B e C: ate 10 metros.",
                   "SINALIZACAO: placa fotoluminescente obrigatoria acima do extintor.",
                   "ACESSO: nunca obstruir com moveis, caixas ou outros objetos.",
                   "QUANTIDADE: calculada pela NBR 12693 conforme area e risco do ambiente."],
             notes=(
                 "A NR-23 e a NBR doze mil seiscentos e noventa e três definem como posicionar os extintores. "
                 "A parte inferior deve estar a no máximo um metro e sessenta do piso. "
                 "Para incêndios da Classe A, a distância máxima a percorrer é quinze metros. "
                 "Para as Classes B e C, dez metros, dada a maior velocidade de propagação. "
                 "A sinalização com placa fotoluminescente é obrigatória acima de cada extintor. "
                 "O acesso deve ser sempre livre — nunca posicione móveis ou caixas na frente do extintor."
             )),
        dict(title="NR-23 — Proteção Contra Incêndios no Trabalho",
             body=["Obrigacoes do empregador:",
                   "• Saidas de emergencia suficientes, sinalizadas e SEMPRE desobstruidas.",
                   "• Extintores adequados ao risco em numero e tipo corretos.",
                   "• BRIGADA DE INCENDIO: funcionarios treinados para resposta inicial.",
                   "• Treinamento anual da brigada com simulado pratico.",
                   "• PPCI aprovado pelo Corpo de Bombeiros Militar.",
                   "• Plano de emergencia e evacuacao com rotas de fuga sinalizadas.",
                   "Fiscalizacao: Corpo de Bombeiros + Ministerio do Trabalho e Emprego."],
             notes=(
                 "A NR-23 trata da proteção contra incêndios no ambiente de trabalho. "
                 "O empregador deve garantir saídas de emergência sempre livres e sinalizadas, "
                 "extintores adequados ao tipo de risco de cada área, "
                 "e manter uma brigada de incêndio formada por funcionários treinados. "
                 "O treinamento da brigada deve ser anual com simulado prático. "
                 "O Plano de Prevenção e Proteção Contra Incêndios deve ser aprovado "
                 "pelo Corpo de Bombeiros Militar de cada estado."
             )),
        dict(title="Resumo — O que Saber sobre Extintores",
             body=["Triangulo do Fogo: Combustivel + Comburente + Calor — remova 1 e extingue.",
                   "Classes: A (solidos) · B (liquidos/gases) · C (eletrico) · D (metais) · K (oleo).",
                   "NUNCA use agua em Classes B, C ou K.",
                   "Tecnica PASS: Puxar pino · Apontar base · Apertar gatilho · Varrer.",
                   "Inspecao mensal + manutencao anual (INMETRO) + hidrostatico a cada 5 anos.",
                   "Distancias: 15 m (Cl.A) / 10 m (Cl.B e C) — acesso livre, sinalizado.",
                   "Apos qualquer uso: recarregar imediatamente."],
             notes=(
                 "Vamos revisar os pontos essenciais sobre extintores. "
                 "O fogo precisa de três elementos e eliminar qualquer um deles o extingue. "
                 "Lembre as cinco classes de incêndio e nunca use água em fogo elétrico, "
                 "de líquidos ou de gordura quente. "
                 "Use sempre a técnica PASS: puxar, apontar para a base, apertar e varrer. "
                 "Mantenha os extintores inspecionados mensalmente e com manutenção anual. "
                 "Respeite as distâncias máximas e mantenha o acesso sempre livre. "
                 "Após qualquer uso, recarregue imediatamente. "
                 "Segurança contra incêndio começa na prevenção."
             )),
    ]

    for s in slides:
        add_slide(prs, s["title"], s["body"], s["notes"], s.get("is_cover", False))
    path = "05_Extintores.pptx"
    prs.save(path)
    print(f"OK: {path} salvo ({len(slides)} slides)")


if __name__ == "__main__":
    print("Gerando apresentacoes PPTX...\n")
    criar_nrs_geral()
    criar_nr10()
    criar_extintores()
    print("\nProntos! Execute: python gerar_videos_MP4.py")
