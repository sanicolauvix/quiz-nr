"""
=============================================================
  GERADOR DE VÍDEO MP4 COM NARRAÇÃO
  A partir dos arquivos .pptx gerados
=============================================================

COMO USAR:
  1. Instale as dependências (veja abaixo)
  2. Coloque este script na mesma pasta dos arquivos .pptx
  3. Execute: python gerar_videos_MP4.py

DEPENDÊNCIAS (execute no terminal):
  pip install python-pptx gTTS moviepy pillow

  No Windows, também é necessário instalar o LibreOffice:
  https://www.libreoffice.org/download/download-libreoffice/

  No Linux/Mac:
  sudo apt install libreoffice  (Linux)
  brew install libreoffice       (Mac)

=============================================================
"""

import os
import sys
import time
import subprocess
import shutil
from pathlib import Path

# ---- Verificar dependências ----
def check_deps():
    missing = []
    try:
        from pptx import Presentation
    except ImportError:
        missing.append("python-pptx")
    try:
        from gtts import gTTS
    except ImportError:
        missing.append("gTTS")
    try:
        import moviepy
    except ImportError:
        missing.append("moviepy")
    try:
        from PIL import Image
    except ImportError:
        missing.append("pillow")
    if missing:
        print(f"❌ Dependências faltando: {', '.join(missing)}")
        print(f"   Execute: pip install {' '.join(missing)}")
        sys.exit(1)
    print("✅ Dependências OK")

check_deps()

from pptx import Presentation
from gtts import gTTS
from moviepy import AudioFileClip
from PIL import Image
import tempfile

def encontrar_libreoffice():
    """Encontra o executável do LibreOffice."""
    candidatos = [
        "soffice", "libreoffice",
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        "/usr/bin/soffice", "/usr/bin/libreoffice",
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
    ]
    for c in candidatos:
        if shutil.which(c) or os.path.exists(c):
            return c
    return None

def pptx_para_imagens(pptx_path, pasta_saida):
    """Converte slides PPTX em imagens PNG usando LibreOffice."""
    soffice = encontrar_libreoffice()
    if not soffice:
        print("❌ LibreOffice não encontrado! Instale em: https://www.libreoffice.org/")
        sys.exit(1)

    print(f"  Convertendo {Path(pptx_path).name} para PDF...")
    pdf_path = str(pasta_saida / "slides.pdf")

    cmd = [soffice, "--headless", "--convert-to", "pdf",
           "--outdir", str(pasta_saida), str(pptx_path)]
    result = subprocess.run(cmd, capture_output=True, text=True)
    if result.returncode != 0:
        print(f"Erro LibreOffice: {result.stderr}")
        sys.exit(1)

    # Encontrar o PDF gerado
    pdf_files = list(pasta_saida.glob("*.pdf"))
    if not pdf_files:
        print("❌ PDF não gerado pelo LibreOffice")
        sys.exit(1)
    pdf_path = pdf_files[0]

    print(f"  Convertendo PDF para imagens...")
    # Usar pdftoppm se disponível, senão usar pdf2image
    try:
        subprocess.run(["pdftoppm", "-v"], capture_output=True, check=True)
        prefix = str(pasta_saida / "slide")
        subprocess.run(
            ["pdftoppm", "-jpeg", "-r", "150", str(pdf_path), prefix],
            check=True, capture_output=True
        )
        imagens = sorted(pasta_saida.glob("slide-*.jpg")) or sorted(pasta_saida.glob("slide*.jpg"))
    except (subprocess.CalledProcessError, FileNotFoundError):
        try:
            import fitz  # pymupdf
            doc = fitz.open(str(pdf_path))
            imagens = []
            for i, page in enumerate(doc):
                mat = fitz.Matrix(150/72, 150/72)
                pix = page.get_pixmap(matrix=mat)
                img_path = pasta_saida / f"slide_{i+1:03d}.jpg"
                pix.save(str(img_path))
                imagens.append(img_path)
            doc.close()
        except ImportError:
            try:
                from pdf2image import convert_from_path
                pages = convert_from_path(str(pdf_path), dpi=150)
                imagens = []
                for i, page in enumerate(pages):
                    img_path = pasta_saida / f"slide_{i+1:03d}.jpg"
                    page.save(str(img_path), "JPEG")
                    imagens.append(img_path)
            except Exception:
                print("❌ Instale pymupdf: pip install pymupdf")
                sys.exit(1)

    return sorted(imagens)

def extrair_notas(pptx_path):
    """Extrai as notas de orador de cada slide."""
    prs = Presentation(pptx_path)
    notas = []
    for slide in prs.slides:
        try:
            texto = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            texto = ""
        notas.append(texto if texto else "Avançar para o próximo slide.")
    return notas

def gerar_audio(texto, caminho_mp3):
    """Gera áudio MP3 a partir de texto usando gTTS."""
    tts = gTTS(text=texto, lang='pt', tld='com.br', slow=False)
    tts.save(caminho_mp3)

def encontrar_ffmpeg():
    """Encontra o executável do ffmpeg."""
    # Tenta via moviepy primeiro
    try:
        from moviepy.config import get_setting
        caminho = get_setting("FFMPEG_BINARY")
        if caminho and os.path.exists(caminho):
            return caminho
    except Exception:
        pass
    # Tenta via imageio-ffmpeg
    try:
        import imageio_ffmpeg
        return imageio_ffmpeg.get_ffmpeg_exe()
    except Exception:
        pass
    # Tenta no PATH
    ffmpeg = shutil.which("ffmpeg")
    if ffmpeg:
        return ffmpeg
    return None

def criar_video(pptx_path, pasta_saida):
    """Processo completo: PPTX → imagens + áudios → MP4 (via ffmpeg direto)."""
    nome_base = Path(pptx_path).stem
    pasta = Path(pasta_saida) / nome_base
    pasta.mkdir(parents=True, exist_ok=True)

    print(f"\n{'='*55}")
    print(f"  Processando: {nome_base}")
    print(f"{'='*55}")

    ffmpeg = encontrar_ffmpeg()
    if not ffmpeg:
        print("❌ ffmpeg não encontrado! Execute: pip install imageio-ffmpeg")
        sys.exit(1)

    # 1. Converter slides para imagens
    print("\n[1/3] Convertendo slides para imagens...")
    imagens = pptx_para_imagens(Path(pptx_path), pasta)
    print(f"      {len(imagens)} slide(s) convertido(s)")

    # 2. Extrair notas, gerar áudios e segmentos de vídeo
    print("\n[2/3] Gerando narrações e segmentos de vídeo...")
    notas = extrair_notas(pptx_path)
    segmentos = []

    for i, (img_path, nota) in enumerate(zip(imagens, notas)):
        slide_n = i + 1
        print(f"      Slide {slide_n}/{len(imagens)}: áudio + vídeo...")

        audio_path = str(pasta / f"audio_{slide_n:03d}.mp3")
        gerar_audio(nota, audio_path)

        # Duração do áudio
        audio_clip = AudioFileClip(audio_path)
        duracao = audio_clip.duration + 0.5
        audio_clip.close()

        # Redimensionar imagem
        img = Image.open(str(img_path))
        img_resized = img.resize((1280, 720), Image.LANCZOS)
        # Garantir dimensões pares (exigido pelo libx264)
        w, h = img_resized.size
        w = w if w % 2 == 0 else w - 1
        h = h if h % 2 == 0 else h - 1
        img_resized = img_resized.resize((w, h))
        img_path_hd = str(pasta / f"slide_{slide_n:03d}_hd.jpg")
        img_resized.save(img_path_hd, "JPEG", quality=90)

        # Criar segmento MP4 com ffmpeg diretamente
        seg_path = str(pasta / f"seg_{slide_n:03d}.mp4")
        cmd = [
            ffmpeg, "-y",
            "-loop", "1", "-framerate", "1", "-i", img_path_hd,
            "-i", audio_path,
            "-c:v", "libx264", "-tune", "stillimage", "-preset", "fast",
            "-c:a", "aac", "-b:a", "128k",
            "-pix_fmt", "yuv420p",
            "-t", str(duracao),
            seg_path
        ]
        result = subprocess.run(cmd, capture_output=True, text=True)
        if result.returncode != 0:
            print(f"Erro ffmpeg no slide {slide_n}: {result.stderr[-300:]}")
            sys.exit(1)
        segmentos.append(seg_path)

    # 3. Concatenar segmentos
    print(f"\n[3/3] Concatenando segmentos em MP4 final...")
    concat_file = str(pasta / "concat.txt")
    with open(concat_file, "w", encoding="utf-8") as f:
        for seg in segmentos:
            f.write(f"file '{seg}'\n")

    mp4_path = str(Path(pasta_saida) / f"{nome_base}.mp4")
    cmd = [
        ffmpeg, "-y",
        "-f", "concat", "-safe", "0",
        "-i", concat_file,
        "-c", "copy",
        mp4_path
    ]
    result = subprocess.run(cmd, capture_output=True, text=True)
    if result.returncode != 0:
        print(f"Erro na concatenação: {result.stderr[-300:]}")
        sys.exit(1)

    tamanho = os.path.getsize(mp4_path) / (1024 * 1024)
    print(f"      ✅ Vídeo gerado: {Path(mp4_path).name} ({tamanho:.1f} MB)")
    return mp4_path

# ============================================================
# EXECUÇÃO PRINCIPAL
# ============================================================
if __name__ == "__main__":
    pasta_script = Path(__file__).parent

    # Encontrar todos os PPTX na mesma pasta
    arquivos_pptx = sorted(pasta_script.glob("*.pptx"))

    if not arquivos_pptx:
        print("❌ Nenhum arquivo .pptx encontrado na pasta!")
        print(f"   Pasta verificada: {pasta_script}")
        sys.exit(1)

    print(f"\n{'='*55}")
    print(f"  GERADOR DE VÍDEOS MP4 COM NARRAÇÃO")
    print(f"{'='*55}")
    print(f"\nArquivos encontrados: {len(arquivos_pptx)}")
    for f in arquivos_pptx:
        print(f"  • {f.name}")

    pasta_saida = pasta_script / "videos_gerados"
    pasta_saida.mkdir(exist_ok=True)
    print(f"\nVídeos serão salvos em: {pasta_saida}\n")

    inicio = time.time()
    videos_gerados = []

    for pptx in arquivos_pptx:
        try:
            mp4 = criar_video(str(pptx), str(pasta_saida))
            videos_gerados.append(mp4)
        except Exception as e:
            print(f"\n❌ Erro ao processar {pptx.name}: {e}")
            import traceback
            traceback.print_exc()

    tempo_total = time.time() - inicio
    print(f"\n{'='*55}")
    print(f"  CONCLUÍDO em {tempo_total:.0f} segundos!")
    print(f"{'='*55}")
    print(f"\n{len(videos_gerados)} vídeo(s) gerado(s):")
    for v in videos_gerados:
        tamanho = os.path.getsize(v) / (1024*1024)
        print(f"  ✅ {Path(v).name} ({tamanho:.1f} MB)")

    print(f"\nPasta de saída: {pasta_saida}")
